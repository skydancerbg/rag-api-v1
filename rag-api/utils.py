# utils.py v1
from pathlib import Path
import json
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract

def read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")

def read_pdf(path: Path) -> str:
    doc = fitz.open(str(path))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    return "\n".join(pages)

def read_docx(path: Path) -> str:
    doc = docx.Document(str(path))
    return "\n".join([p.text for p in doc.paragraphs])

def read_pptx(path: Path) -> str:
    pres = Presentation(str(path))
    texts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def read_html(path: Path) -> str:
    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text(separator="\n")

def read_json(path: Path) -> str:
    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(data, ensure_ascii=False)
    except Exception:
        return path.read_text(encoding="utf-8", errors="ignore")

def ocr_image(path: Path) -> str:
    try:
        img = Image.open(str(path))
        return pytesseract.image_to_string(img)
    except Exception:
        return ""

def extract_text(path):
    p = Path(path)
    ext = p.suffix.lower()
    if ext in [".txt", ".md"]:
        return read_text_file(p)
    if ext == ".pdf":
        return read_pdf(p)
    if ext == ".docx":
        return read_docx(p)
    if ext == ".pptx":
        return read_pptx(p)
    if ext == ".html" or ext == ".htm":
        return read_html(p)
    if ext == ".json":
        return read_json(p)
    if ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        return ocr_image(p)
    return ""
